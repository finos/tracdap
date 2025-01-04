#  Licensed to the Fintech Open Source Foundation (FINOS) under one or
#  more contributor license agreements. See the NOTICE file distributed
#  with this work for additional information regarding copyright ownership.
#  FINOS licenses this file to you under the Apache License, Version 2.0
#  (the "License"); you may not use this file except in compliance with the
#  License. You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

import typing as _tp

import tracdap.rt.api as trac

import pptx
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

import tutorial.schemas as schemas


class PowerpointReport(trac.TracModel):

    def define_parameters(self) -> _tp.Dict[str, trac.ModelParameter]:

        return trac.define_parameters()

    def define_inputs(self) -> _tp.Dict[str, trac.ModelInputSchema]:

        quarterly_sales_schema = trac.load_schema(schemas, "profit_by_region.csv")
        quarterly_sales = trac.define_input(quarterly_sales_schema, label="Quarterly sales data")
        report_template = trac.define_input(trac.CommonFileTypes.POWERPOINT, label="Quarterly sales report template")

        return {
            "quarterly_sales": quarterly_sales,
            "report_template": report_template }

    def define_outputs(self) -> _tp.Dict[str, trac.ModelOutputSchema]:

        sales_report = trac.define_output(trac.CommonFileTypes.POWERPOINT, label="Quarterly sales report")

        return { "sales_report": sales_report }

    def run_model(self, ctx: trac.TracContext):

        quarterly_sales = ctx.get_pandas_table("quarterly_sales")

        with ctx.get_file_stream("report_template") as report_template:
            presentation = pptx.Presentation(report_template)

        slide = presentation.slides.add_slide(presentation.slide_layouts[5])

        title_frame = slide.shapes[0].text_frame
        title_frame.text = 'Profit by Region Report'

        # define chart data ---------------------
        chart_data = CategoryChartData()
        chart_data.categories = quarterly_sales["region"]
        chart_data.add_series('Gross Profit', quarterly_sales["gross_profit"])

        # add chart to slide --------------------
        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
        slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)

        with ctx.put_file_stream("sales_report") as sales_report:
            presentation.save(sales_report)


if __name__ == "__main__":
    import tracdap.rt.launch as launch
    launch.launch_model(PowerpointReport, "config/file_io_powerpoint.yaml", "config/sys_config.yaml")
